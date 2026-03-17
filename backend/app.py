from flask import Flask, request, jsonify
from flask_cors import CORS
from PyPDF2 import PdfReader
from pptx import Presentation
from sentence_transformers import SentenceTransformer
import faiss
import numpy as np

app = Flask(__name__)
CORS(app)

# Load model once globally
model = SentenceTransformer('all-MiniLM-L6-v2')

# In-memory storage
stored_chunks = []
stored_index = None

def chunk_text(text, chunk_size=500):
    return [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]

@app.route('/upload', methods=['POST'])
def upload_file():
    global stored_chunks, stored_index

    if 'file' not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file = request.files['file']
    if file.filename == "":
        return jsonify({"error": "Empty file name"}), 400

    text = ""
    try:
        if file.filename.lower().endswith('.pdf'):
            reader = PdfReader(file)
            for page in reader.pages:
                text += page.extract_text() or ""
        elif file.filename.lower().endswith('.pptx'):
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        else:
            return jsonify({"error": "Unsupported file type"}), 400

        if not text.strip():
            return jsonify({"error": "No text found in file"}), 400

        chunks = chunk_text(text)
        embeddings = model.encode(chunks).astype('float32') # Fix: FAISS needs float32

        dimension = embeddings.shape[1]
        index = faiss.IndexFlatL2(dimension)
        index.add(embeddings)

        stored_chunks = chunks
        stored_index = index

        return jsonify({"message": f"Successfully indexed {len(chunks)} chunks"})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/ask', methods=['POST'])
def ask_question():
    global stored_chunks, stored_index

    if stored_index is None:
        return jsonify({"answer": "Please upload a file first!"}), 400

    data = request.get_json()
    question = data.get('question')

    if not question:
        return jsonify({"answer": "Please provide a question"}), 400

    try:
        q_embedding = model.encode([question]).astype('float32')
        # Search for top 3 matches
        D, I = stored_index.search(q_embedding, k=min(3, len(stored_chunks)))
        
        results = [stored_chunks[i] for i in I[0] if i != -1]
        answer = " ".join(results) if results else "No matches found."

        return jsonify({"answer": answer})
    except Exception as e:
        return jsonify({"answer": f"Error: {str(e)}"}), 500

if __name__ == '__main__':
    app.run(debug=True, port=5000)
